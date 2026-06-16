"""
Document text extraction service.

Supports PDF (PyMuPDF + pdfplumber fallback), PPTX, DOCX, and TXT.
Also performs basic text cleaning (whitespace normalization, noise removal)
while preserving heading-like lines.
"""
import os
import re
from pathlib import Path

import fitz  # PyMuPDF
import pdfplumber
from pptx import Presentation
import docx


class TextExtractionError(Exception):
    pass


def extract_text_from_pdf(file_path: str) -> str:
    """Extract text from a PDF using PyMuPDF, falling back to pdfplumber per-page."""
    pages_text: list[str] = []
    try:
        doc = fitz.open(file_path)
        for page in doc:
            text = page.get_text("text")
            if not text.strip():
                text = ""
            pages_text.append(text)
        doc.close()
    except Exception:
        pages_text = []

    # Fallback for pages that produced no text (e.g. scanned/odd encoding)
    if not pages_text or all(not p.strip() for p in pages_text):
        pages_text = []
        with pdfplumber.open(file_path) as pdf:
            for page in pdf.pages:
                pages_text.append(page.extract_text() or "")

    return "\n\n".join(
        f"[PAGE {i + 1}]\n{text}" for i, text in enumerate(pages_text)
    )


def extract_text_from_pptx(file_path: str) -> str:
    """Extract text from a PPTX, slide by slide, preserving titles."""
    prs = Presentation(file_path)
    slides_text = []
    for idx, slide in enumerate(prs.slides, start=1):
        parts = [f"[SLIDE {idx}]"]
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    line = "".join(run.text for run in para.runs).strip()
                    if line:
                        parts.append(line)
        slides_text.append("\n".join(parts))
    return "\n\n".join(slides_text)


def extract_text_from_docx(file_path: str) -> str:
    """Extract text from a DOCX, preserving heading styles as markdown-like headers."""
    document = docx.Document(file_path)
    lines = []
    for para in document.paragraphs:
        text = para.text.strip()
        if not text:
            continue
        style = (para.style.name or "").lower()
        if "heading" in style:
            level = "".join(filter(str.isdigit, style)) or "1"
            lines.append(f"{'#' * min(int(level), 6)} {text}")
        else:
            lines.append(text)

    # Tables
    for table in document.tables:
        for row in table.rows:
            cells = [cell.text.strip() for cell in row.cells]
            if any(cells):
                lines.append(" | ".join(cells))

    return "\n".join(lines)


def extract_text_from_txt(file_path: str) -> str:
    with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
        return f.read()


EXTRACTORS = {
    ".pdf": extract_text_from_pdf,
    ".pptx": extract_text_from_pptx,
    ".ppt": extract_text_from_pptx,
    ".docx": extract_text_from_docx,
    ".txt": extract_text_from_txt,
}


def extract_text(file_path: str) -> str:
    """Dispatch to the correct extractor based on file extension."""
    ext = Path(file_path).suffix.lower()
    extractor = EXTRACTORS.get(ext)
    if not extractor:
        raise TextExtractionError(f"Unsupported file extension: {ext}")
    try:
        raw_text = extractor(file_path)
    except Exception as exc:
        raise TextExtractionError(f"Failed to extract text from {file_path}: {exc}") from exc

    return clean_text(raw_text)


def clean_text(text: str) -> str:
    """
    Clean extracted text:
    - Normalize line endings
    - Collapse repeated blank lines / whitespace
    - Strip trailing whitespace per line
    - Preserve heading markers ([PAGE], [SLIDE], #...) on their own lines
    """
    text = text.replace("\r\n", "\n").replace("\r", "\n")

    lines = []
    for raw_line in text.split("\n"):
        line = re.sub(r"[ \t]+", " ", raw_line).strip()
        lines.append(line)

    # Collapse 3+ consecutive blank lines to 1
    cleaned_lines: list[str] = []
    blank_streak = 0
    for line in lines:
        if line == "":
            blank_streak += 1
            if blank_streak > 1:
                continue
        else:
            blank_streak = 0
        cleaned_lines.append(line)

    cleaned = "\n".join(cleaned_lines).strip()
    # Remove duplicate spaces left from joins
    cleaned = re.sub(r" {2,}", " ", cleaned)
    return cleaned


def save_extracted_text(text: str, output_path: str) -> None:
    os.makedirs(os.path.dirname(output_path), exist_ok=True)
    with open(output_path, "w", encoding="utf-8") as f:
        f.write(text)
